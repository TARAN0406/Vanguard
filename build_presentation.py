from pptx import Presentation

prs = Presentation(r'C:\Users\navte\Downloads\Finspark_Hackathon_Template.pptx')

# We'll do slide-specific replacements based on slide index (0-based)
replacements = {
    0: {
        "Your Team Name :": "Your Team Name : Vanguard",
        "Your team bio :": "Your team bio : We are a team of passionate developers tackling internal threat vectors.",
        "Date :": "Date : Finspark Hackathon 2026"
    },
    1: {
        "Mention the problem statement selected by your team and explain why you chose to solve it.": "Problem Statement 1: Insider Threat Detection.\nWe chose this because internal threats cause significant financial losses, and standard SIEMs lack AI-driven contextual zero-trust enforcement."
    },
    2: {
        "List any assumptions, required access, datasets, tools, APIs, environments or inputs needed for your solution to work.": "Pre-requisites:\n- Node.js environment\n- Web browser\n- Built-in SQLite database (no setup required)\n- Gemini API Key (optional for LLM features)"
    },
    3: {
        "Mention the technologies, frameworks, platforms, libraries, datasets or references used to build your solution.": "Tech Stack:\n- Frontend: React 19, Tailwind CSS, Vite\n- Backend: Node.js, Express, SQLite\n- AI/ML: TensorFlow.js, Isolation Forest, Google Gemini\n- Visuals: Recharts"
    },
    4: {
        "Add any supporting documents, user flows, process notes, research links, logic flows or references that explain your solution better.": "Supporting Documents:\n- docs/ARCHITECTURE.md (System Design)\n- docs/SETUP.md (Installation Guide)\n- docs/Vanguard_Project_Summary.md"
    },
    5: {
        "Mention the technologies, frameworks, platforms, libraries, datasets or references used to build your solution.": "Key Differentiators:\n1. 3-Layer AI Ensemble (Z-Score + Isolation Forest + Deep Autoencoder)\n2. Actionable '1-Click Freeze' remediation\n3. Decoupled architecture with zero-trust design\n\nAdoption Plan:\nEasy drop-in replacement for existing internal portals with lightweight DB footprint."
    },
    6: {
        "Add your GitHub link along with any diagrams, screenshots or architecture references that support your prototype.": "GitHub: https://github.com/NaVteJSingH18/Vanguard\n\n(See architecture diagram in the repo's docs folder)"
    },
    7: {
        "Explain the future potential of your solution, including expansion, real-world use cases and long-term applicability.": "Business Potential:\n- Can be sold as a white-label SIEM to mid-sized banks.\n- Easily scales to multi-tenant environments.\n- Significantly reduces data exfiltration risks."
    },
    8: {
        "Highlight what is innovative, original or different about your approach compared to existing solutions.": "Uniqueness:\nWe don't just 'alert'. We use an AI ensemble to accurately calculate risk, and leverage Gemini to summarize the threat context in plain English for rapid SOC action."
    },
    9: {
        "Show how users will interact with the solution and why it is simple, clear and usable.": "User Experience:\n- Built for SOC Analysts\n- Dark mode glassmorphism UI\n- Real-time threat telemetry charts\n- Clean, distraction-free environment"
    },
    10: {
        "Explain how the solution can scale across branches, users, systems, transaction volumes or banking environments.": "Scalability:\n- The React frontend is completely decoupled from the Node backend.\n- SQLite can be instantly swapped for PostgreSQL to handle millions of transactions.\n- Stateless JWT auth enables horizontal scaling."
    },
    11: {
        "Explain how practical the solution is to implement, run, maintain and improve over time.": "Deployment:\n- Run 'npm install' and 'npm start'\n- Fully containerizable via Docker\n- Zero mandatory cloud dependencies for the core ML engine (runs locally in Node)"
    },
    12: {
        "Describe how the solution handles cybersecurity risks, data protection, access control, compliance and safe implementation.": "Security:\n- Strict Role-Based Access Control (RBAC)\n- Immutable Audit Logs via SQLite\n- Zero-Trust validation on every endpoint\n- JWT authentication"
    },
    13: {
        "Add a clear architecture diagram showing how the solution works across components, systems, data flows and users.": "Architecture:\n[Frontend UI] <--> [Node.js / Express API]\n[API] <--> [AI Scoring Engine (TF.js / Isolation Forest)]\n[API] <--> [SQLite Database]\n[API] <--> [Gemini LLM for Threat Intel]"
    },
    14: {
        "Add prototype screenshots, demo video link, GitHub link and any visuals that prove the solution is functional.": "Links & Media:\n- GitHub: https://github.com/NaVteJSingH18/Vanguard\n- Screenshots available in 'assets/screenshots/' directory of the repository."
    },
    15: {
        "Add team member names, roles, contact details and any closing note.": "Team Vanguard\n- Navtej Singh (Lead Developer)\n\nThank you for reviewing our Zero-Trust Intelligence platform!"
    }
}

for i, slide in enumerate(prs.slides):
    if i in replacements:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for old_text, new_text in replacements[i].items():
                    if old_text in shape.text:
                        shape.text = shape.text.replace(old_text, new_text)

prs.save(r'C:\Users\navte\.gemini\antigravity\scratch\Vanguard\Vanguard_Presentation.pptx')
print("Successfully generated Vanguard_Presentation.pptx")
